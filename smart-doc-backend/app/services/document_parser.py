"""文档解析服务 - 支持 PDF/Word/PPT/TXT/图片"""
import os
import tempfile
from werkzeug.utils import secure_filename

def parse_pptx(file_path):
    """解析 PPTX 文件，提取所有页文字"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                text += f"第{slide_idx + 1}页:\n" + "\n".join(slide_text) + "\n\n"
        
        if not text.strip():
            text = "【PPT 中未提取到文字】"
        return text.strip()
    except Exception as e:
        return f"【PPT 解析失败: {str(e)}】"


def parse_docx(file_path):
    """解析 Word 文件"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if not text.strip():
            text = "【Word 文档中未提取到文字】"
        return text
    except Exception as e:
        return f"【Word 解析失败: {str(e)}】"


def parse_pdf(file_path):
    """解析 PDF 文件"""
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if not text.strip():
            text = "【PDF 中未提取到文字，可能是扫描件】"
        return text.strip()
    except Exception as e:
        return f"【PDF 解析失败: {str(e)}】"


def parse_txt(file_path):
    """解析 TXT 文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='gbk') as f:
            return f.read()


def parse_image(file_path):
    """图片文件返回 None，走 OCR 流程"""
    return None


def parse_document(file_path, filename):
    """根据文件扩展名选择解析器"""
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        return parse_pdf(file_path)
    elif ext == 'docx':
        return parse_docx(file_path)
    elif ext == 'pptx':
        return parse_pptx(file_path)
    elif ext == 'txt':
        return parse_txt(file_path)
    elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'gif']:
        return None  # 图片，走 OCR
    else:
        return f"【不支持的文件格式: {ext}】"